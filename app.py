import streamlit as st
import pandas as pd
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase.pdfmetrics import stringWidth
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import zipfile
import os
from PIL import Image

st.title("🎨 手作り英単語カードメーカー")

# --- フォントの設定 ---
font_path = "comicbd.ttf" 
if os.path.exists(font_path):
    pdfmetrics.registerFont(TTFont('ComicSans', font_path))
    target_font = 'ComicSans'
    pptx_font_name = 'Comic Sans MS'
else:
    target_font = 'Helvetica-Bold'
    pptx_font_name = 'Arial'

col1, col2 = st.columns(2)
with col1:
    csv_file = st.file_uploader("1. 単語リスト(CSV)", type=['csv'])
with col2:
    zip_file = st.file_uploader("2. 画像まとめ(ZIP)", type=['zip'])

output_type = st.radio("3. 出力形式を選択", ["PDF", "PowerPoint (PPTX)"])

if csv_file and zip_file:
    df = pd.read_csv(csv_file, header=None)
    words = [str(w) for w in df[0].tolist()]
    z = zipfile.ZipFile(zip_file)
    file_list = z.namelist()

    if st.button(f"{output_type} を作成する"):
        buf = io.BytesIO()
        
        if output_type == "PDF":
            c = canvas.Canvas(buf, pagesize=landscape(A4))
            width, height = landscape(A4)
            for word in words:
                max_f = height * 0.4
                curr_f = max_f
                while curr_f > 10:
                    if stringWidth(word, target_font, curr_f) <= width * 0.9: break
                    curr_f -= 5
                c.setFont(target_font, curr_f)
                c.drawCentredString(width / 2, (height / 2) - (curr_f / 3), word)
                c.showPage()
                found = None
                for ext in ['.jpg', '.jpeg', '.png', '.JPG', '.PNG']:
                    target = f"{word}{ext}"
                    for f in file_list:
                        if f.endswith(f"/{target}") or f == target:
                            found = f; break
                    if found: break
                if found:
                    img = Image.open(io.BytesIO(z.read(found)))
                    c.drawInlineImage(img, 0, 0, width=width, height=height, preserveAspectRatio=False)
                c.showPage()
            c.save()
            file_ext, mime = "pdf", "application/pdf"

        else: 
            # --- PowerPoint 2010 徹底対策版 ---
            prs = Presentation() # スライドサイズ設定を削除（デフォルトの4:3を使用）
            
            for word in words:
                # 表面：文字
                # レイアウト0（タイトルスライド）が最も安定します
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                # 標準のタイトル枠を使用
                title = slide.shapes.title
                title.text = word
                # フォント設定
                for paragraph in title.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(80)
                        run.font.name = pptx_font_name
                        run.font.bold = True
                # 副題（Subtitle）があれば削除してエラーを回避
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        sp = shape.element
                        sp.getparent().remove(sp)
                
                # 裏面：画像
                slide = prs.slides.add_slide(prs.slide_layouts[6]) # 空白スライド
                found = None
                for ext in ['.jpg', '.jpeg', '.png', '.JPG', '.PNG']:
                    target = f"{word}{ext}"
                    for f in file_list:
                        if f.endswith(f"/{target}") or f == target:
                            found = f; break
                    if found: break
                if found:
                    img_data = io.BytesIO(z.read(found))
                    # サイズ指定を「スライドいっぱい」ではなく、少し余裕を持たせる
                    slide.shapes.add_picture(img_data, 0, 0, width=prs.slide_width)
            
            prs.save(buf)
            file_ext, mime = "pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        st.success(f"{output_type} が完成しました！")
        st.download_button(label="保存して2010で試す", data=buf.getvalue(), file_name=f"English_Cards_2010.{file_ext}", mime=mime)
